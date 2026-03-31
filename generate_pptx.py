"""
Generate PayShield cybersecurity hackathon presentation.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ──────────────────────────────────────────────────────────
BG       = RGBColor(0x1a, 0x1a, 0x2e)   # dark navy
TITLE_BG = RGBColor(0xe7, 0x4c, 0x3c)   # red accent
BLUE     = RGBColor(0x29, 0x80, 0xb9)   # info blue
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT    = RGBColor(0xCC, 0xCC, 0xCC)   # light grey for body
CARD_BG  = RGBColor(0x16, 0x21, 0x3e)   # slightly lighter than BG
CARD2_BG = RGBColor(0x0f, 0x3c, 0x56)   # blue-tinted card

# ── Slide dimensions (widescreen 16:9) ──────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK_LAYOUT = prs.slide_layouts[6]   # completely blank


# ────────────────────────────────────────────────────────────────────────────
# Helper functions
# ────────────────────────────────────────────────────────────────────────────

def set_bg(slide, color=BG):
    """Fill slide background with a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=None):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, wrap=True, italic=False):
    """Add a text box with a single paragraph."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_title_bar(slide, title_text, subtitle_text=None):
    """Add the red title bar at the top of every slide."""
    bar = add_rect(slide, Inches(0), Inches(0), W, Inches(1.1), TITLE_BG)

    # Title
    add_text_box(slide, title_text,
                 Inches(0.3), Inches(0.08), Inches(12.5), Inches(0.7),
                 font_size=32, bold=True, color=WHITE)

    # Optional subtitle inside bar
    if subtitle_text:
        add_text_box(slide, subtitle_text,
                     Inches(0.3), Inches(0.72), Inches(12.5), Inches(0.35),
                     font_size=14, color=RGBColor(0xFF, 0xDD, 0xDD))


def bullet_frame(slide, bullets, left, top, width, height,
                 font_size=17, color=LIGHT, indent=0.2):
    """Add a multi-bullet text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for bullet in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        p.level = 0
        run = p.add_run()
        run.text = "• " + bullet
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return txBox


def add_card(slide, left, top, width, height, bg_color,
             lines, font_sizes=None, colors=None, border_color=None):
    """Add a card (rounded rect replaced by rect) with stacked text lines."""
    add_rect(slide, left, top, width, height, bg_color,
             line_color=border_color or TITLE_BG, line_width=Pt(1.5))
    if font_sizes is None:
        font_sizes = [16] * len(lines)
    if colors is None:
        colors = [WHITE] * len(lines)
    y = top + Inches(0.12)
    for line, fs, col in zip(lines, font_sizes, colors):
        tb = add_text_box(slide, line,
                          left + Inches(0.12), y,
                          width - Inches(0.24), Inches(0.4),
                          font_size=fs, color=col, wrap=True)
        y += Inches(0.38)


# ────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Contexte PayShield
# ────────────────────────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s1)
add_title_bar(s1, "Contexte PayShield", "Gestion des Vulnérabilités Assistée par IA")

bullets = [
    "PayShield : fintech béninoise, 280 employés, 15 000 tx/jour, 1,2M utilisateurs",
    "Suite à une alerte CERT : scan complet Nessus + OpenVAS → 35 vulnérabilités sur 10 systèmes",
    "Risque : compromission de la base de paiement, violation PCI-DSS, impact direct utilisateurs",
    "Problème : le score CVSS brut ne permet pas de prioriser la remédiation dans ce contexte métier",
    "Solution : un algorithme de scoring contextuel SRC enrichissant le CVSS avec 3 facteurs PayShield",
]
bullet_frame(s1, bullets, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.8), font_size=18)


# ────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Notre formule SRC
# ────────────────────────────────────────────────────────────────────────────
s2 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s2)
add_title_bar(s2, "Notre Formule SRC")

# Formula box
formula_box = add_rect(s2, Inches(1), Inches(1.25), Inches(11.33), Inches(0.65), CARD_BG,
                       line_color=TITLE_BG, line_width=Pt(2))
add_text_box(s2, "SRC(v)  =  CVSS × W_exposition × W_criticité × W_données",
             Inches(1.1), Inches(1.3), Inches(11.1), Inches(0.55),
             font_size=20, bold=True, color=TITLE_BG, align=PP_ALIGN.CENTER)

# Table header
def trow(slide, y, col1, col2, col3, header=False):
    row_h = Inches(0.42)
    bg = TITLE_BG if header else CARD_BG
    tc = WHITE if header else LIGHT
    for x, w, txt in [
        (Inches(0.3),  Inches(2.3), col1),
        (Inches(2.65), Inches(4.5), col2),
        (Inches(7.2),  Inches(5.8), col3),
    ]:
        add_rect(slide, x, y, w, row_h, bg,
                 line_color=RGBColor(0x44,0x44,0x66), line_width=Pt(0.5))
        add_text_box(slide, txt, x + Inches(0.08), y + Inches(0.04),
                     w - Inches(0.16), row_h - Inches(0.08),
                     font_size=14 if not header else 15,
                     bold=header, color=tc)

trow(s2, Inches(2.05), "Facteur", "Valeurs", "Justification métier", header=True)
trow(s2, Inches(2.50), "Exposition",
     "Internet ×2.0  /  DMZ ×1.5  /  Interne ×1.0",
     "Une faille accessible depuis Internet = risque d'exploitation immédiat")
trow(s2, Inches(2.95), "Criticité",
     "Critique ×2.0  /  Élevée ×1.5  /  Moyenne ×1.0  /  Faible ×0.5",
     "La BDD paiement (Critique) ne peut pas être traitée comme dev-env")
trow(s2, Inches(3.40), "Données",
     "Sensible ×1.5  /  Interne ×1.0  /  Public ×0.5",
     "KYC et transactions paiement = données à protéger en priorité")

# Note
add_text_box(s2, "SRC_norm = SRC / max(SRC) × 10   →   score normalisé 0–10",
             Inches(0.5), Inches(4.05), Inches(12.3), Inches(0.45),
             font_size=15, color=BLUE, bold=True, align=PP_ALIGN.CENTER)


# ────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Top 10 — Résultats
# ────────────────────────────────────────────────────────────────────────────
s3 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s3)
add_title_bar(s3, "Top 10 — Résultats SRC V1")

add_text_box(s3, "Les 3 vulnérabilités prioritaires",
             Inches(0.5), Inches(1.18), Inches(12.3), Inches(0.4),
             font_size=18, bold=True, color=BLUE)

cards = [
    {
        "rank": "#1",
        "id": "V01 — OpenSSL (Heartbleed)",
        "ip": "IP: 10.0.0.10",
        "cvss": "CVSS 10.0",
        "src": "SRC 10.0",
        "tags": "Internet / Critique / Sensible",
        "desc": "Fuite mémoire SSL — clés TLS et tokens de session exposés",
    },
    {
        "rank": "#2",
        "id": "V02 — Apache RCE",
        "ip": "IP: 10.0.0.10",
        "cvss": "CVSS 9.8",
        "src": "SRC 9.80",
        "tags": "Internet / Critique / Sensible",
        "desc": "Path traversal sans auth — exécution de commandes distantes",
    },
    {
        "rank": "#3",
        "id": "V34 — Kubernetes API",
        "ip": "IP: 10.0.0.10",
        "cvss": "CVSS 9.8",
        "src": "SRC 9.80",
        "tags": "Internet / Critique / Sensible",
        "desc": "API non authentifiée — prise de contrôle du cluster de paiement",
    },
]

card_w = Inches(4.1)
card_h = Inches(4.5)
gap    = Inches(0.26)
start_x = Inches(0.3)

for i, c in enumerate(cards):
    cx = start_x + i * (card_w + gap)
    cy = Inches(1.65)
    add_rect(s3, cx, cy, card_w, card_h, CARD_BG,
             line_color=TITLE_BG, line_width=Pt(2))
    # rank badge
    add_rect(s3, cx, cy, Inches(0.7), Inches(0.52), TITLE_BG)
    add_text_box(s3, c["rank"], cx + Inches(0.04), cy + Inches(0.04),
                 Inches(0.62), Inches(0.44), font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    y_off = cy + Inches(0.6)
    for txt, fs, col, bd in [
        (c["id"],   17, WHITE,       True),
        (c["ip"],   14, LIGHT,       False),
        (c["cvss"] + "   " + c["src"], 20, TITLE_BG, True),
        (c["tags"], 13, BLUE,        False),
        (c["desc"], 14, LIGHT,       False),
    ]:
        add_text_box(s3, txt, cx + Inches(0.1), y_off,
                     card_w - Inches(0.2), Inches(0.55),
                     font_size=fs, color=col, bold=bd, wrap=True)
        y_off += Inches(0.6) if fs != 14 else Inches(0.52)


# ────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Problématique + Analyse
# ────────────────────────────────────────────────────────────────────────────
s4 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s4)
add_title_bar(s4, "Pourquoi le CVSS Seul ne Suffit Pas")

sections = [
    ("Le problème", [
        "V05 (PostgreSQL, CVSS 10) et V35 (Confluence, CVSS 10) — même score, risques très différents",
        "CVSS ignore : exposition réseau, criticité système, sensibilité des données",
    ]),
    ("Expérience A — Impact des coefficients", [
        "Config Base vs Config Alt : Internet ×2.0→×3.0, Sensible ×1.5→×2.0, Interne ×1.0→×0.5",
        "Résultat : même Top 10, mais scores internes s'effondrent — V05 passe de 5.0 → 1.67",
    ]),
    ("Anomalie identifiée", [
        "V05 (PostgreSQL BDD paiement, CVSS 10) classé 11e — hors Top 10",
        "Cause : malus exposition Interne (×1.0) pénalise la BDD de paiement",
        "Correction : intégrer l'EPSS → SRC_v2 = SRC × (1 + epss)",
    ]),
]

y = Inches(1.3)
for title, items in sections:
    add_rect(s4, Inches(0.3), y, Inches(12.7), Inches(0.38), CARD2_BG)
    add_text_box(s4, title, Inches(0.42), y + Inches(0.04),
                 Inches(12.5), Inches(0.32), font_size=15, bold=True, color=BLUE)
    y += Inches(0.38)
    bullet_frame(s4, items, Inches(0.5), y, Inches(12.3),
                 Inches(len(items) * 0.5 + 0.1), font_size=15)
    y += Inches(len(items) * 0.5 + 0.25)


# ────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — SRC V3 — EPSS + KEV
# ────────────────────────────────────────────────────────────────────────────
s5 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s5)
add_title_bar(s5, "SRC V3 — EPSS + KEV")

# Formula
add_rect(s5, Inches(1), Inches(1.25), Inches(11.33), Inches(0.65), CARD_BG,
         line_color=TITLE_BG, line_width=Pt(2))
add_text_box(s5, "SRC_v3  =  SRC × (1 + epss) × (1 + 1.5 × kev)",
             Inches(1.1), Inches(1.30), Inches(11.1), Inches(0.55),
             font_size=20, bold=True, color=TITLE_BG, align=PP_ALIGN.CENTER)

# Two sub-section cards
for x, text in [
    (Inches(0.3),  "(1 + epss) : probabilité d'exploitation dans 30 jours.\nEx: EPSS 0.97 → ×1.97"),
    (Inches(6.75), "(1 + 1.5 × kev) : exploit actif CISA KEV.\nKEV=1 → ×2.5  /  KEV=0 → inchangé"),
]:
    add_rect(s5, x, Inches(2.05), Inches(6.3), Inches(0.9), CARD2_BG,
             line_color=BLUE, line_width=Pt(1))
    add_text_box(s5, text, x + Inches(0.1), Inches(2.1),
                 Inches(6.1), Inches(0.8), font_size=14, color=LIGHT, wrap=True)

# Comparison table
trow2_y = [Inches(3.1), Inches(3.54), Inches(3.98), Inches(4.42)]

def trow2(slide, y, c1, c2, c3, header=False):
    row_h = Inches(0.42)
    bg = TITLE_BG if header else CARD_BG
    tc = WHITE
    for x, w, txt in [
        (Inches(0.3),  Inches(2.5), c1),
        (Inches(2.85), Inches(3.5), c2),
        (Inches(6.4),  Inches(6.6), c3),
    ]:
        add_rect(slide, x, y, w, row_h, bg,
                 line_color=RGBColor(0x44,0x44,0x66), line_width=Pt(0.5))
        add_text_box(slide, txt, x + Inches(0.08), y + Inches(0.04),
                     w - Inches(0.16), row_h - Inches(0.08),
                     font_size=14, bold=header, color=tc)

trow2(s5, trow2_y[0], "Version", "Formule", "Impact sur V35 Confluence", header=True)
trow2(s5, trow2_y[1], "V1 Base",    "SRC = CVSS × W",                   "Rang 26e — ignoré")
trow2(s5, trow2_y[2], "V2 EPSS",    "SRC × (1+epss)",                   "Rang amélioré — EPSS 0.94 pris en compte")
trow2(s5, trow2_y[3], "V3 EPSS+KEV","SRC × (1+epss) × (1+1.5×kev)",    "Remonte significativement — KEV=1 × 2.5")

# Limit note
add_text_box(s5, "Limite : EPSS et KEV ne couvrent pas toutes les CVE — données manquantes possibles",
             Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.4),
             font_size=13, color=BLUE, italic=True)


# ────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Plan J+7 pour PayShield
# ────────────────────────────────────────────────────────────────────────────
s6 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s6)
add_title_bar(s6, "Plan de Remédiation J+7",
              "3 actions immédiates — sans ces actions, PayShield reste exposé")

actions = [
    {
        "num": "Action 1",
        "title": "Patcher OpenSSL sur pay-api (10.0.0.10)",
        "resp": "Responsable : Équipe DevOps",
        "impact": "Impact si non fait : Fuite de clés TLS et tokens — 1,2M comptes compromis",
    },
    {
        "num": "Action 2",
        "title": "Définir mot de passe PostgreSQL sur db-prod (10.0.0.15)",
        "resp": "Responsable : DBA + RSSI",
        "impact": "Impact si non fait : Accès root BDD paiement — vol de toutes les transactions",
    },
    {
        "num": "Action 3",
        "title": "Isoler l'API Kubernetes du réseau public (10.0.0.10)",
        "resp": "Responsable : Équipe Cloud",
        "impact": "Impact si non fait : Prise de contrôle du cluster — déploiement de code malveillant",
    },
]

card_w = Inches(4.0)
card_h = Inches(3.9)
gap    = Inches(0.3)
start_x = Inches(0.37)

for i, a in enumerate(actions):
    cx = start_x + i * (card_w + gap)
    cy = Inches(1.35)
    add_rect(s6, cx, cy, card_w, card_h, CARD_BG,
             line_color=TITLE_BG, line_width=Pt(2))
    # action number badge
    add_rect(s6, cx, cy, card_w, Inches(0.5), TITLE_BG)
    add_text_box(s6, a["num"], cx + Inches(0.1), cy + Inches(0.06),
                 card_w - Inches(0.2), Inches(0.4),
                 font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(s6, a["title"],
                 cx + Inches(0.12), cy + Inches(0.62),
                 card_w - Inches(0.24), Inches(0.9),
                 font_size=15, bold=True, color=WHITE, wrap=True)

    add_text_box(s6, a["resp"],
                 cx + Inches(0.12), cy + Inches(1.6),
                 card_w - Inches(0.24), Inches(0.4),
                 font_size=13, color=BLUE, wrap=True)

    # impact warning box
    add_rect(s6, cx + Inches(0.12), cy + Inches(2.1),
             card_w - Inches(0.24), Inches(1.6),
             RGBColor(0x3d, 0x10, 0x10),
             line_color=TITLE_BG, line_width=Pt(1))
    add_text_box(s6, a["impact"],
                 cx + Inches(0.2), cy + Inches(2.18),
                 card_w - Inches(0.4), Inches(1.42),
                 font_size=13, color=TITLE_BG, bold=True, wrap=True)

# Footer note
add_text_box(s6,
             "J+30 : Keycloak, pfSense, fail2ban   |   J+90 : WAF, segmentation réseau, politique IAM",
             Inches(0.5), Inches(5.45), Inches(12.3), Inches(0.45),
             font_size=13, color=LIGHT, italic=True, align=PP_ALIGN.CENTER)


# ────────────────────────────────────────────────────────────────────────────
# Save
# ────────────────────────────────────────────────────────────────────────────
out = "/home/modlvkindns/Documents/Myworkdir/Projet EPITECH/TP/hackaton-cyber-ai/presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
